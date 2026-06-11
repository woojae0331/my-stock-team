"""report-pptx 렌더러 — 구조화 spec(JSON)을 셀사이드 리서치 덱 톤의 PPTX로 그립니다.

렌더링만 담당합니다. 리포트 .md를 spec으로 정리하는 일은 SKILL.md 지침에 따라 Claude가 합니다.

사용법:
    python build_pptx.py --spec reports/삼성전자.spec.json --out reports/삼성전자.pptx

슬라이드 순서(고정):
    표지 → 종목 개요 → 재무 요약 → 가격/추세 → 뉴스·심리 → 리스크 → 한 줄 종합
"""

from __future__ import annotations

import argparse
import json
import re

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── 디자인 시스템 ────────────────────────────────────────────
FONT_KR = "맑은 고딕"  # 한글 폰트 고정 — 글자 깨짐 방지

KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)  # 포인트색(절제해서 사용)
INK = RGBColor(0x16, 0x1A, 0x24)        # 제목·강조 (near-black)
NAVY = RGBColor(0x23, 0x2A, 0x3D)       # 표 헤더·차트 시리즈
BODY = RGBColor(0x3D, 0x42, 0x4D)       # 본문
CAPTION = RGBColor(0x8A, 0x90, 0x9C)    # 출처·캡션
RULE = RGBColor(0xDD, 0xE0, 0xE6)       # 가는 구분선
CARD = RGBColor(0xF5, 0xF6, 0xF8)       # 카드/밴딩 배경
BAND = RGBColor(0xF7, 0xF8, 0xFA)       # 표 교차 음영
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CHART_PALETTE = [NAVY, KB_YELLOW, RGBColor(0xA7, 0xAD, 0xBA)]

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MX = Inches(0.7)                 # 좌우 여백
CONTENT_W = SLIDE_W - 2 * MX

# 세로 기준선
Y_EYEBROW = Inches(0.52)
Y_TITLE = Inches(0.82)
Y_RULE = Inches(1.46)
Y_BODY = Inches(1.74)
Y_FOOT_RULE = Inches(6.92)
Y_FOOT = Inches(7.02)

TOTAL_SLIDES = 7
DISCLAIMER = "학습·참고용 분석이며 투자 권유가 아닙니다. 매수·매도 판단은 투자자 본인의 몫입니다."
MAX_TABLE_ROWS = 11  # 헤더 포함. 초과 시 잘라내고 안내행.


# ── 저수준 헬퍼 ──────────────────────────────────────────────
def set_font(run, size=13, bold=False, color=BODY, name=FONT_KR, spacing=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    if spacing is not None:
        rPr.set("spc", str(int(spacing)))  # 자간(1/100 pt) — 케른닝 느낌
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def rect(slide, left, top, width, height, fill, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def hline(slide, left, top, width, color=RULE, weight=Pt(1)):
    return rect(slide, left, top, width, weight, color)


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, first, runs, space_after=8, line_spacing=1.12, align=PP_ALIGN.LEFT):
    """runs = [(text, {font kwargs}), ...]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = 0
    p.line_spacing = line_spacing
    for text, kw in runs:
        r = p.add_run()
        r.text = text
        set_font(r, **kw)
    return p


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    return slide


def _num(s):
    """문자열에서 숫자 추출(차트용). %·비파싱·결측은 None."""
    s = str(s)
    if "%" in s:
        return None
    m = re.search(r"-?\d[\d,]*\.?\d*", s)
    return float(m.group().replace(",", "")) if m else None


# ── 공통 프레임(헤더/푸터) ───────────────────────────────────
def chrome(slide, *, section, title, running, page):
    # 좌상단 섹션 넘버(옐로우) + 섹션 라벨
    tf = textbox(slide, MX, Y_EYEBROW, CONTENT_W * 0.7, Inches(0.3))
    para(
        tf, True,
        [(f"{section:02d}", dict(size=11, bold=True, color=KB_YELLOW, spacing=60)),
         ("   ", dict(size=11)),
         ("SECTION", dict(size=9, bold=True, color=CAPTION, spacing=180))],
        space_after=0,
    )
    # 우상단 러닝 식별자
    tfr = textbox(slide, SLIDE_W - MX - Inches(4.5), Y_EYEBROW, Inches(4.5), Inches(0.3))
    para(tfr, True, [(running, dict(size=10, color=CAPTION, spacing=40))], space_after=0, align=PP_ALIGN.RIGHT)

    # 제목
    tft = textbox(slide, MX, Y_TITLE, CONTENT_W, Inches(0.6))
    para(tft, True, [(title, dict(size=24, bold=True, color=INK))], space_after=0)

    # 제목 구분선 + 좌측 옐로우 탭
    hline(slide, MX, Y_RULE, CONTENT_W, RULE, Pt(1))
    rect(slide, MX, Y_RULE - Pt(1), Inches(0.85), Pt(3), KB_YELLOW)

    # 푸터
    hline(slide, MX, Y_FOOT_RULE, CONTENT_W, RULE, Pt(0.75))
    tff = textbox(slide, MX, Y_FOOT, CONTENT_W * 0.8, Inches(0.32))
    para(tff, True, [(DISCLAIMER, dict(size=8.5, color=CAPTION))], space_after=0)
    tfp = textbox(slide, SLIDE_W - MX - Inches(1.6), Y_FOOT, Inches(1.6), Inches(0.32))
    para(tfp, True, [(f"{page:02d} / {TOTAL_SLIDES:02d}", dict(size=9, bold=True, color=NAVY))],
         space_after=0, align=PP_ALIGN.RIGHT)


def bullets(slide, items, *, top=Y_BODY, left=MX, width=CONTENT_W, height=Inches(5.0), size=13.5, gap=11):
    tf = textbox(slide, left, top, width, height)
    for i, item in enumerate(items):
        para(
            tf, i == 0,
            [("—  ", dict(size=size, bold=True, color=KB_YELLOW)),
             (str(item), dict(size=size, color=BODY))],
            space_after=gap, line_spacing=1.18,
        )


def caption(slide, text, top, left=MX, width=CONTENT_W):
    tf = textbox(slide, left, top, width, Inches(0.3))
    txt = text if text.strip().startswith("(") else f"출처/기준일: {text}"
    para(tf, True, [(txt, dict(size=9, color=CAPTION))], space_after=0)


# ── KPI 카드 ─────────────────────────────────────────────────
def kpi_cards(slide, cards, *, top=Y_BODY, height=Inches(1.5)):
    n = len(cards)
    gap = Inches(0.28)
    cw = Emu(int((CONTENT_W - gap * (n - 1)) / n))
    for i, c in enumerate(cards):
        left = Emu(int(MX + i * (cw + gap)))
        rect(slide, left, top, cw, height, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(slide, left, top, Inches(0.5), Pt(3), KB_YELLOW)  # 상단 옐로우 틱
        tf = textbox(slide, Emu(int(left + Inches(0.22))), Emu(int(top + Inches(0.18))),
                     Emu(int(cw - Inches(0.44))), Emu(int(height - Inches(0.3))))
        para(tf, True, [(str(c.get("label", "")), dict(size=10.5, bold=True, color=CAPTION, spacing=30))],
             space_after=6)
        para(tf, False, [(str(c.get("value", "")), dict(size=24, bold=True, color=INK))], space_after=4)
        if c.get("source"):
            para(tf, False, [(str(c["source"]), dict(size=8, color=CAPTION))], space_after=0, line_spacing=1.05)


# ── 표 ───────────────────────────────────────────────────────
def table(slide, columns, rows, *, left=MX, top=Y_BODY, width=CONTENT_W, max_rows=MAX_TABLE_ROWS):
    rows = [list(map(str, r)) for r in rows]
    body_cap = max_rows - 1
    if len(rows) > body_cap:
        omitted = len(rows) - (body_cap - 1)
        rows = rows[: body_cap - 1]
        rows.append([f"… 외 {omitted}개 행 생략 …"] + [""] * (len(columns) - 1))

    n_rows, n_cols = len(rows) + 1, len(columns)
    hdr_h, body_h = Inches(0.46), Inches(0.4)
    total_h = Emu(int(hdr_h + body_h * (n_rows - 1)))
    gfx = slide.shapes.add_table(n_rows, n_cols, left, top, width, total_h)
    tbl = gfx.table

    # 기본 표 스타일 제거(테두리/밴딩 없는 스타일) → 직접 색칠
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = tbl._tbl.makeelement(qn("a:tblPr"), {})
        tbl._tbl.insert(0, tblPr)
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is None:
        sid = tblPr.makeelement(qn("a:tableStyleId"), {})
        tblPr.append(sid)
    sid.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid

    tbl.rows[0].height = hdr_h
    for r in range(1, n_rows):
        tbl.rows[r].height = body_h

    def fill_cell(cell, text, *, header=False, banded=False, col=0):
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if header else (BAND if banded else WHITE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.14)
        cell.margin_right = Inches(0.14)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        tfc = cell.text_frame
        tfc.word_wrap = True
        p = tfc.paragraphs[0]
        right = col > 0  # 첫 열 외에는 우측 정렬(수치 컬럼)
        p.alignment = PP_ALIGN.RIGHT if right else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        set_font(r, size=11.5, bold=header, color=(WHITE if header else BODY))

    for c, col in enumerate(columns):
        fill_cell(tbl.cell(0, c), str(col), header=True, col=c)
    for ri, row in enumerate(rows, start=1):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            fill_cell(tbl.cell(ri, c), val, banded=(ri % 2 == 0), col=c)

    bottom = Emu(int(top + total_h))
    # 헤더 아래 옐로우 액센트 + 표 하단 가는 선
    rect(slide, left, Emu(int(top + hdr_h)) - Pt(2), width, Pt(2), KB_YELLOW)
    hline(slide, left, bottom, width, RULE, Pt(0.75))
    return bottom


# ── 차트(네이티브, 편집 가능) ────────────────────────────────
def column_chart(slide, categories, series, *, left, top, width, height, unit=""):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.name = FONT_KR
    ch.font.size = Pt(10)

    ch.has_legend = len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(10)
        ch.legend.font.name = FONT_KR

    for i, s in enumerate(ch.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = CHART_PALETTE[i % len(CHART_PALETTE)]
        s.format.line.fill.background()

    plot = ch.plots[0]
    plot.gap_width = 80
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = "#,##0.0"
    dl.number_format_is_linked = False
    dl.font.size = Pt(9)
    dl.font.name = FONT_KR
    dl.font.color.rgb = BODY
    try:
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass

    va = ch.value_axis
    va.has_major_gridlines = False
    va.has_minor_gridlines = False
    try:
        va.visible = False
    except Exception:
        pass
    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(10)
    ca.tick_labels.font.name = FONT_KR
    try:
        ca.format.line.color.rgb = RULE
    except Exception:
        pass
    return gf


def _financial_series(tbl):
    """재무 표에서 차트용 시리즈를 안전하게 추출. 파싱 불가 시 None."""
    if not tbl:
        return None
    cols, rows = tbl.get("columns", []), tbl.get("rows", [])
    cats = [str(c) for c in cols[1:]]
    if len(cats) < 2:
        return None
    series = []
    for row in rows:
        vals = [_num(v) for v in row[1:]]
        if vals and all(v is not None for v in vals) and len(vals) == len(cats):
            series.append((str(row[0]), vals))
    return (cats, series) if len(series) >= 1 else None


# ── 슬라이드 ─────────────────────────────────────────────────
def slide_cover(prs, cover):
    slide = blank(prs)
    name = cover.get("stock_name", "종목")
    code = cover.get("stock_code", "")

    # 상단 가는 옐로우 룰
    rect(slide, 0, 0, SLIDE_W, Inches(0.16), KB_YELLOW)
    # 좌측 여백 기준선
    left = Inches(1.1)

    tf = textbox(slide, left, Inches(1.5), Inches(8), Inches(0.4))
    para(tf, True, [("EQUITY RESEARCH  ·  학습·참고용", dict(size=12, bold=True, color=CAPTION, spacing=160))],
         space_after=0)

    tft = textbox(slide, left, Inches(2.35), Inches(11), Inches(1.6))
    para(tft, True, [(name, dict(size=48, bold=True, color=INK))], space_after=6)
    # 티커 칩
    if code:
        chip = rect(slide, left + Inches(0.02), Inches(3.55), Inches(1.9), Inches(0.42),
                    NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
        ctf = chip.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = f"KRX {code}"
        set_font(cr, size=12, bold=True, color=WHITE, spacing=40)

    para(textbox(slide, left, Inches(4.25), Inches(11), Inches(0.5)), True,
         [(cover.get("subtitle", "종목 리서치 리포트"), dict(size=16, color=BODY))], space_after=0)

    # 메타 룰 + 작성일/통화
    hline(slide, left, Inches(5.2), Inches(10.8), RULE, Pt(1))
    mtf = textbox(slide, left, Inches(5.4), Inches(11), Inches(0.4))
    para(mtf, True,
         [("작성일  ", dict(size=11, bold=True, color=CAPTION, spacing=40)),
          (f"{cover.get('date', '')}      ", dict(size=11, color=BODY)),
          ("통화  ", dict(size=11, bold=True, color=CAPTION, spacing=40)),
          ("KRW", dict(size=11, color=BODY))],
         space_after=0)

    # 하단 면책
    ftf = textbox(slide, left, Inches(6.85), Inches(11.2), Inches(0.4))
    para(ftf, True, [(DISCLAIMER, dict(size=9, color=CAPTION))], space_after=0)
    return slide


def slide_overview(prs, data, running, page):
    slide = blank(prs)
    chrome(slide, section=1, title="종목 개요", running=running, page=page)
    top = Y_BODY
    if data.get("kpis"):
        kpi_cards(slide, data["kpis"][:4], top=Y_BODY, height=Inches(1.5))
        top = Y_BODY + Inches(1.85)
    bullets(slide, data.get("bullets", []), top=top, height=Emu(int(Y_FOOT_RULE - top - Inches(0.1))))
    return slide


def slide_financials(prs, data, running, page):
    slide = blank(prs)
    chrome(slide, section=2, title="재무 요약", running=running, page=page)
    tbl = data.get("table")
    parsed = _financial_series(tbl)

    if tbl and parsed:
        # 좌: 표 / 우: 차트
        half = Emu(int((CONTENT_W - Inches(0.5)) / 2))
        tbl_bottom = table(slide, tbl["columns"], tbl["rows"], left=MX, top=Y_BODY, width=half)
        if tbl.get("source"):
            caption(slide, tbl["source"], tbl_bottom + Inches(0.12), left=MX, width=half)
        cats, series = parsed
        chart_left = Emu(int(MX + half + Inches(0.5)))
        column_chart(slide, cats, series, left=chart_left, top=Y_BODY - Inches(0.05),
                     width=half, height=Inches(3.0))
        cap2 = textbox(slide, chart_left, Y_BODY + Inches(3.05), half, Inches(0.3))
        para(cap2, True, [("단위: 조 원 · 출처는 좌측 표 기준", dict(size=8.5, color=CAPTION))],
             space_after=0, align=PP_ALIGN.CENTER)
        body_top = Y_BODY + Inches(3.55)
    elif tbl:
        tbl_bottom = table(slide, tbl["columns"], tbl["rows"], top=Y_BODY)
        if tbl.get("source"):
            caption(slide, tbl["source"], tbl_bottom + Inches(0.12))
        body_top = tbl_bottom + Inches(0.55)
    else:
        body_top = Y_BODY

    if data.get("bullets"):
        bullets(slide, data["bullets"], top=body_top,
                height=Emu(int(Y_FOOT_RULE - body_top - Inches(0.1))), size=12.5, gap=8)
    return slide


def slide_price_trend(prs, data, running, page):
    slide = blank(prs)
    chrome(slide, section=3, title="가격 / 추세", running=running, page=page)
    image = data.get("image")
    tbl = data.get("table")
    if image:
        try:
            slide.shapes.add_picture(image, MX, Y_BODY, width=Emu(int(CONTENT_W * 0.55)))
        except Exception:
            pass
        if tbl:
            right = Emu(int(MX + CONTENT_W * 0.58))
            w = Emu(int(CONTENT_W * 0.42))
            tb = table(slide, tbl["columns"], tbl["rows"], left=right, top=Y_BODY, width=w)
            if tbl.get("source"):
                caption(slide, tbl["source"], tb + Inches(0.12), left=right, width=w)
        body_top = Y_BODY + Inches(3.6)
    elif tbl:
        # 좌: KPI 강조 / 우: 표  — 표만 있을 땐 표를 좌측, 코멘트 우측
        half = Emu(int((CONTENT_W - Inches(0.5)) / 2))
        tb = table(slide, tbl["columns"], tbl["rows"], left=MX, top=Y_BODY, width=half)
        if tbl.get("source"):
            caption(slide, tbl["source"], tb + Inches(0.12), left=MX, width=half)
        if data.get("bullets"):
            bullets(slide, data["bullets"], top=Y_BODY,
                    left=Emu(int(MX + half + Inches(0.5))), width=half,
                    height=Inches(4.5), size=12.5, gap=9)
        return slide
    else:
        body_top = Y_BODY
    if data.get("bullets"):
        bullets(slide, data["bullets"], top=body_top,
                height=Emu(int(Y_FOOT_RULE - body_top - Inches(0.1))), size=12.5)
    return slide


def slide_news(prs, data, running, page):
    slide = blank(prs)
    chrome(slide, section=4, title="뉴스 · 심리", running=running, page=page)
    sentiment = data.get("sentiment")
    top = Y_BODY
    if sentiment:
        # 심리 배너
        rect(slide, MX, Y_BODY, CONTENT_W, Inches(0.62), CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(slide, MX, Y_BODY, Inches(0.12), Inches(0.62), KB_YELLOW)
        tf = textbox(slide, MX + Inches(0.32), Y_BODY, CONTENT_W - Inches(0.5), Inches(0.62),
                     anchor=MSO_ANCHOR.MIDDLE)
        para(tf, True,
             [("시장 심리   ", dict(size=12, bold=True, color=CAPTION, spacing=60)),
              (str(sentiment), dict(size=14, bold=True, color=INK))],
             space_after=0)
        top = Y_BODY + Inches(0.95)
    bullets(slide, data.get("bullets", []), top=top,
            height=Emu(int(Y_FOOT_RULE - top - Inches(0.1))), size=13.5)
    return slide


def slide_risk(prs, data, running, page):
    slide = blank(prs)
    chrome(slide, section=5, title="리스크", running=running, page=page)
    items = data.get("bullets", [])
    # 리스크는 번호 카드 형태로 — 좌측 옐로우 인덱스
    top = Y_BODY
    avail = Emu(int(Y_FOOT_RULE - Y_BODY - Inches(0.15)))
    h = Emu(int((avail - Inches(0.25) * (len(items) - 1)) / max(len(items), 1))) if items else avail
    h = min(h, Inches(1.4))
    for i, it in enumerate(items):
        y = Emu(int(top + i * (h + Inches(0.25))))
        rect(slide, MX, y, CONTENT_W, h, CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(slide, MX, y, Inches(0.12), h, KB_YELLOW)
        idx = textbox(slide, MX + Inches(0.3), y, Inches(0.8), h, anchor=MSO_ANCHOR.MIDDLE)
        para(idx, True, [(f"{i + 1:02d}", dict(size=22, bold=True, color=RGBColor(0xC9, 0xCE, 0xD8)))], space_after=0)
        body = textbox(slide, MX + Inches(1.15), y, CONTENT_W - Inches(1.5), h, anchor=MSO_ANCHOR.MIDDLE)
        text = re.sub(r"^\s*\d+\.\s*", "", str(it))  # 앞 번호 제거(카드가 번호 부여)
        para(body, True, [(text, dict(size=12.5, color=BODY))], space_after=0, line_spacing=1.15)
    return slide


def slide_summary(prs, data, running, page):
    slide = blank(prs)
    chrome(slide, section=6, title="한 줄 종합", running=running, page=page)
    rect(slide, MX, Inches(2.7), CONTENT_W, Inches(1.9), CARD, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, MX, Inches(2.7), Inches(0.16), Inches(1.9), KB_YELLOW)
    tf = textbox(slide, MX + Inches(0.6), Inches(2.9), CONTENT_W - Inches(1.1), Inches(1.5),
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, True, [("핵심 종합", dict(size=11, bold=True, color=CAPTION, spacing=120))], space_after=10)
    para(tf, False, [(data.get("text", ""), dict(size=17, bold=True, color=INK))],
         space_after=0, line_spacing=1.3)
    return slide


def build(spec, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover = spec.get("cover", {})
    running = f"{cover.get('stock_name', '')} · {cover.get('stock_code', '')}".strip(" ·")

    slide_cover(prs, cover)
    slide_overview(prs, spec.get("overview", {}), running, 2)
    slide_financials(prs, spec.get("financials", {}), running, 3)
    slide_price_trend(prs, spec.get("price_trend", {}), running, 4)
    slide_news(prs, spec.get("news", {}), running, 5)
    slide_risk(prs, spec.get("risk", {}), running, 6)
    slide_summary(prs, spec.get("summary", {}), running, 7)

    prs.save(out_path)
    print(f"PPTX 생성 완료: {out_path} (슬라이드 {len(prs.slides)}장)")


def main():
    ap = argparse.ArgumentParser(description="구조화 spec(JSON) → 리서치 덱 PPTX")
    ap.add_argument("--spec", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()
    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)
    build(spec, args.out)


if __name__ == "__main__":
    main()
